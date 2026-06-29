"""Tests for PPTXConverter. Fixtures are generated with python-pptx at runtime."""

from pptx import Presentation

from pydown.converters.pptx import PPTXConverter


def _add_slide(prs, title=None, body=None, notes=None):
    # Layout 1 is "Title and Content": a title placeholder plus a body placeholder.
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    if title is not None:
        slide.shapes.title.text = title
    if body is not None:
        slide.placeholders[1].text = body
    if notes is not None:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def test_title_and_body(tmp_path):
    prs = Presentation()
    _add_slide(prs, title="My Title", body="Some body text")
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    md = PPTXConverter().convert(str(path))

    assert "## Slide 1" in md
    assert "My Title" in md
    assert "Some body text" in md


def test_multiple_slides(tmp_path):
    prs = Presentation()
    _add_slide(prs, title="One")
    _add_slide(prs, title="Two")
    path = tmp_path / "multi.pptx"
    prs.save(str(path))

    md = PPTXConverter().convert(str(path))

    assert "## Slide 1" in md
    assert "## Slide 2" in md
    assert md.index("## Slide 1") < md.index("## Slide 2")


def test_speaker_notes(tmp_path):
    prs = Presentation()
    _add_slide(prs, title="Talk", notes="Remember to smile")
    path = tmp_path / "notes.pptx"
    prs.save(str(path))

    md = PPTXConverter().convert(str(path))

    assert "> Notes:" in md
    assert "> Remember to smile" in md


def test_empty_slide(tmp_path):
    prs = Presentation()
    # Layout 6 is blank — no placeholders, no text.
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "blank.pptx"
    prs.save(str(path))

    md = PPTXConverter().convert(str(path))

    # Still emits the slide header, no crash.
    assert "## Slide 1" in md
