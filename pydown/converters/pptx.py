"""PPTX -> Markdown converter, backed by python-pptx.

Each slide becomes a ``## Slide N`` section: the title shape first, then the text
of every other shape, and finally any speaker notes as a ``> Notes:`` blockquote.
"""

from __future__ import annotations

from pptx import Presentation

from .base import BaseConverter


def _blockquote(text):
    """Render ``text`` as a Markdown blockquote, prefixing every line."""
    lines = text.splitlines() or [""]
    return "\n".join(f"> {line}" if line else ">" for line in lines)


class PPTXConverter(BaseConverter):
    """Convert a .pptx presentation into Markdown, one section per slide."""

    extensions = (".pptx",)

    def _convert(self, file_path: str) -> str:
        presentation = Presentation(file_path)
        sections = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts = [f"## Slide {slide_number}"]

            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title = title_shape.text.strip()
                if title:
                    parts.append(title)

            # Body text from every non-title shape that holds text.
            for shape in slide.shapes:
                if shape is title_shape or not shape.has_text_frame:
                    continue
                text = shape.text.strip()
                if text:
                    parts.append(text)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append("> Notes:\n" + _blockquote(notes))

            sections.append("\n\n".join(parts))

        return "\n\n".join(sections)
